"""Build the PROJECT_OVERVIEW summary deck.

A 4-part briefing that a newcomer can absorb in ~40 minutes:
  Part 1: StreamingLLM paper (Xiao et al., ICLR 2024)
  Part 2: What we did for Task 1
  Part 3: GPTCache open-source project
  Part 4: What we did for Task 2

Palette and helpers mirror `task1-presentation/build_streaming_llm_deck.py`
so the two decks feel like siblings.

Author: Nissim Brami (nissimbrami@post.bgu.ac.il)
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


HERE = Path(__file__).resolve().parent
OUT = HERE / "PROJECT_OVERVIEW.pptx"


NAVY = RGBColor(0x0B, 0x1F, 0x3A)
BLUE = RGBColor(0x2E, 0x86, 0xAB)
ACCENT = RGBColor(0xE6, 0x7E, 0x22)
TEXT = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x66, 0x66, 0x66)
GREEN_OK = RGBColor(0x27, 0xAE, 0x60)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _title(slide, text, top_in=0.35):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(top_in),
                                  Inches(12.33), Inches(0.9))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = NAVY
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(top_in + 0.95),
                                      Inches(12.83), Inches(top_in + 0.95))
    line.line.color.rgb = ACCENT
    line.line.width = Pt(2.25)


def _footer(slide, page_num, total):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.05),
                                  Inches(12.33), Inches(0.3))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Nissim Brami \u00b7 Caching in LLMs \u00b7 BGU"
    r.font.name = "Calibri"
    r.font.size = Pt(9)
    r.font.color.rgb = MUTED
    tb2 = slide.shapes.add_textbox(Inches(12.3), Inches(7.05),
                                   Inches(0.5), Inches(0.3))
    tf2 = tb2.text_frame
    tf2.margin_left = 0
    tf2.margin_top = 0
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = f"{page_num} / {total}"
    r2.font.name = "Calibri"
    r2.font.size = Pt(9)
    r2.font.color.rgb = MUTED


def _bullets(slide, items, top_in=1.55, left_in=0.6, width_in=12.13,
             height_in=5.2, font_size=18):
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in),
                                  Inches(width_in), Inches(height_in))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        if isinstance(item, tuple):
            head, body = item
            r = p.add_run()
            r.text = f"\u2022 {head} \u2014 "
            r.font.name = "Calibri"
            r.font.size = Pt(font_size)
            r.font.bold = True
            r.font.color.rgb = NAVY
            r2 = p.add_run()
            r2.text = body
            r2.font.name = "Calibri"
            r2.font.size = Pt(font_size)
            r2.font.color.rgb = TEXT
        else:
            r = p.add_run()
            r.text = f"\u2022 {item}"
            r.font.name = "Calibri"
            r.font.size = Pt(font_size)
            r.font.color.rgb = TEXT


def _section_banner(slide, part, subtitle):
    # Full-bleed navy banner behind title
    from pptx.enum.shapes import MSO_SHAPE
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0), Inches(2.6),
                                  Inches(13.33), Inches(2.3))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.85),
                                  Inches(12.33), Inches(0.7))
    tf = tb.text_frame
    tf.margin_left = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = part
    r.font.name = "Calibri"
    r.font.size = Pt(20)
    r.font.color.rgb = ACCENT
    r.font.bold = True

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.5),
                                   Inches(12.33), Inches(1.2))
    tf2 = tb2.text_frame
    tf2.margin_left = 0
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.name = "Calibri"
    r2.font.size = Pt(32)
    r2.font.bold = True
    r2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slides = []  # list of (kind, payload)

    # 1. Title
    slides.append(("title", None))
    # 2. Agenda
    slides.append(("agenda", None))

    # --- Part 1: StreamingLLM paper ---
    slides.append(("part", ("PART 1",
                            "StreamingLLM \u2014 the paper")))
    slides.append(("bullets", ("The problem: streaming LLMs break",
                               [("Setting",
                                 "Chatbots, agents, long-running assistants \u2014 "
                                 "input length T grows unboundedly."),
                                ("KV cache",
                                 "Every token stores a Key and Value vector. "
                                 "Memory grows O(T); attention grows O(T\u00b2)."),
                                ("Three broken baselines (Llama-2-13B, 20k tokens)",
                                 "Dense: 5641 PPL (OOM). "
                                 "Window (last 2048): 5158 PPL. "
                                 "Window+recompute: 5.43 PPL but O(TL\u00b2) slow."),
                                ("So what?",
                                 "None of these is deployable at 4 million tokens. "
                                 "Something fundamental is missing.")])))
    slides.append(("bullets", ("The observation: attention sinks",
                               [("Attention maps show a shocking pattern",
                                 "Layers 3+ dump 30\u201360% of attention mass onto "
                                 "the first 4 tokens \u2014 no matter what those tokens are."),
                                ("Why it happens",
                                 "Softmax must sum to 1. When no other token is "
                                 "relevant, the model needs a \"no-op\" bucket. "
                                 "Initial tokens are visible to everyone \u2014 they win."),
                                ("Line-break clincher",
                                 "Replace token 0 with \"\\n\": PPL 5.60 vs 5.40. "
                                 "Position matters far more than content."),
                                ("Two-word fix",
                                 "Keep 4 initial sinks + last L rolling tokens. "
                                 "That is StreamingLLM.")])))
    slides.append(("bullets", ("Result: perplexity flat past 4M tokens",
                               [("Correctness",
                                 "PG-19 book test: PPL stays flat from 4k \u2192 4M tokens. "
                                 "Baselines diverge at ~2\u00d7 pretraining length."),
                                ("Speed",
                                 "Up to 22.2\u00d7 faster than sliding-window + recompute "
                                 "(Llama-2-13B, A6000). Memory is O(L), not O(T)."),
                                ("Cost of the fix",
                                 "4 extra KV entries. That is the entire runtime cost."),
                                ("Adopted in production",
                                 "HuggingFace SinkCache, TensorRT-LLM \"streaming\" mode, "
                                 "vLLM sliding-window variants.")])))

    # --- Part 2: Task 1 ---
    slides.append(("part", ("PART 2",
                            "Task 1 \u2014 what we built")))
    slides.append(("bullets", ("Task 1 assignment & deliverables",
                               [("The assignment",
                                 "Pick a caching-in-LLMs paper. Give a 30-min class "
                                 "lecture. Convince Prof. Einziger you understand it."),
                                ("Why StreamingLLM",
                                 "Perfect fit: the KV cache is the classic "
                                 "capacity-limited cache; the paper is a policy paper."),
                                ("What we produced",
                                 "46-slide 16:9 deck, 8000-word speaker notes, "
                                 "runnable PyTorch demo, unit tests, reference decks "
                                 "for continuity with prior BGU students."),
                                ("How we frame it",
                                 "Cache-terminology-first: sinks = admission policy, "
                                 "rolling window = LRU-style eviction.")])))
    slides.append(("bullets", ("Deck structure (12 sections)",
                               [("Basics",
                                 "Slides 1\u201310: what a KV cache is, why streams break it, "
                                 "course vocabulary (LRU/LFU/admission/eviction)."),
                                ("The paper",
                                 "Slides 11\u201325: attention-sink discovery, softmax argument, "
                                 "the 4+L policy, complexity analysis."),
                                ("Results & baselines",
                                 "Slides 26\u201335: PG-19 numbers, speedup, streaming QA."),
                                ("Critique & extensions",
                                 "Slides 36\u201346: limits (no memory of dropped tokens), "
                                 "what I'd change, connection to Task 2 (cost-aware eviction).")])))

    # --- Part 3: GPTCache ---
    slides.append(("part", ("PART 3",
                            "GPTCache \u2014 the open source")))
    slides.append(("bullets", ("What GPTCache is",
                               [("The idea",
                                 "Semantic response cache in front of the LLM API. "
                                 "\"What's the capital of France?\" and \"France's capital?\" "
                                 "get the same cached answer."),
                                ("Origins",
                                 "Zilliz (Milvus vector DB team), open-sourced mid-2023. "
                                 "Apache 2.0, ~12k GitHub stars, Python."),
                                ("Why it matters",
                                 "LLM APIs cost real money. Caching hits by meaning, "
                                 "not by string, is the natural next step after HTTP caching."),
                                ("Three pluggable stores",
                                 "CacheBase (SQL / SQLite / MySQL), "
                                 "VectorBase (FAISS / Milvus / Chroma), "
                                 "EvictionBase (which entries to drop when full).")])))
    slides.append(("bullets", ("The existing eviction policies",
                               [("LRU (Least Recently Used)",
                                 "Kick out whatever hasn't been touched in longest. "
                                 "Ignores frequency, size, cost."),
                                ("LFU (Least Frequently Used)",
                                 "Kick out whichever was used least often. "
                                 "Old popular entries stick; new entries starve."),
                                ("FIFO / RR",
                                 "First-in-first-out / random. Simple, weak."),
                                ("What none of them do",
                                 "Consider that a cache miss on a $0.03 GPT-4 call "
                                 "costs 30\u00d7 more than a miss on a $0.001 embedding. "
                                 "\u2192 our opening.")])))

    # --- Part 4: Task 2 ---
    slides.append(("part", ("PART 4",
                            "Task 2 \u2014 GDSF cost-aware plugin")))
    slides.append(("bullets", ("What we contribute",
                               [("The gap",
                                 "GPTCache assumes all entries cost the same to recompute. "
                                 "They don't."),
                                ("The policy",
                                 "GDSF (Cao & Irani 1997, Cherkasova 1998): "
                                 "Priority = L + freq(i)^\u03b1 \u00b7 cost(i)^\u03b2 / size(i)"),
                                ("Adaptation",
                                 "cost(i) = the recorded $ per API call for that entry. "
                                 "L = clock counter to prevent starvation."),
                                ("What we ship",
                                 "PR-ready branch fork of GPTCache: 1060-line plugin, "
                                 "259 unit tests, benchmark harness, 3600-run study, "
                                 "7-page ACM sigconf paper.")])))
    slides.append(("bullets", ("Results (3600 runs, paired-t + Bonferroni + BCa bootstrap)",
                               [("Cost-Weighted Hit Rate on Zipf-0.9 workloads",
                                 "+25.7% vs LRU, +32.3% vs LFU, +91.0% vs FIFO, "
                                 "+18.8% vs Random."),
                                ("Uniform-cost sanity check",
                                 "\u00b10.005% vs baselines \u2014 policy correctly "
                                 "collapses to LFU-like behaviour when there's "
                                 "nothing to optimise."),
                                ("Overhead",
                                 "-0.037% throughput vs LRU \u2014 statistically indistinguishable. "
                                 "Priority computation is O(1) per access."),
                                ("Statistics",
                                 "All wins Bonferroni-adjusted p<0.001 with 95% BCa CIs "
                                 "computed over 10000 bootstrap resamples (seed 20260721).")])))

    # --- Bridge / recap ---
    slides.append(("bullets", ("How the two tasks connect",
                               [("Same setting",
                                 "Both tasks study a bounded cache serving an LLM "
                                 "workload. The question is only \"what do you evict?\""),
                                ("Task 1 (KV cache)",
                                 "Position-based eviction \u2014 keep 4 sinks + rolling window. "
                                 "Domain: attention arithmetic."),
                                ("Task 2 (response cache)",
                                 "Cost-based eviction \u2014 keep the entries that are "
                                 "most expensive to regenerate. Domain: semantic hits."),
                                ("Bridge",
                                 "Both replace naive FIFO/LRU with a smarter policy "
                                 "that respects the actual value of what's in the cache.")])))

    # --- Thanks ---
    slides.append(("thanks", None))

    total = len(slides)

    for i, (kind, payload) in enumerate(slides, start=1):
        slide = _blank(prs)

        if kind == "title":
            # Title slide
            tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.4),
                                          Inches(12.13), Inches(1.4))
            tf = tb.text_frame
            tf.margin_left = 0
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = "Caching in LLMs \u2014 Project Overview"
            r.font.name = "Calibri"
            r.font.size = Pt(44)
            r.font.bold = True
            r.font.color.rgb = NAVY

            tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(3.8),
                                           Inches(12.13), Inches(0.8))
            tf2 = tb2.text_frame
            p2 = tf2.paragraphs[0]
            r2 = p2.add_run()
            r2.text = ("Four 10-minute reads: StreamingLLM \u00b7 Task 1 \u00b7 "
                       "GPTCache \u00b7 Task 2")
            r2.font.name = "Calibri"
            r2.font.size = Pt(22)
            r2.font.color.rgb = BLUE

            tb3 = slide.shapes.add_textbox(Inches(0.6), Inches(5.6),
                                           Inches(12.13), Inches(0.5))
            tf3 = tb3.text_frame
            p3 = tf3.paragraphs[0]
            r3 = p3.add_run()
            r3.text = ("Nissim Brami \u00b7 Prof. Gil Einziger \u00b7 "
                       "Ben-Gurion University \u00b7 2026")
            r3.font.name = "Calibri"
            r3.font.size = Pt(14)
            r3.font.color.rgb = MUTED

            line = slide.shapes.add_connector(1, Inches(0.6), Inches(2.35),
                                              Inches(12.73), Inches(2.35))
            line.line.color.rgb = ACCENT
            line.line.width = Pt(3)

        elif kind == "agenda":
            _title(slide, "Agenda")
            _bullets(slide, [
                ("Part 1",
                 "The StreamingLLM paper \u2014 the problem, the observation, "
                 "the fix, the numbers."),
                ("Part 2",
                 "What we built for Task 1 \u2014 46-slide deck, notes, demo code."),
                ("Part 3",
                 "GPTCache \u2014 origins, architecture, existing eviction policies."),
                ("Part 4",
                 "What we built for Task 2 \u2014 GDSF plugin, 3600-run benchmark, "
                 "paper."),
                ("Bridge",
                 "How the two tasks connect through cache-policy design."),
            ], font_size=20)
            _footer(slide, i, total)

        elif kind == "part":
            part, subtitle = payload
            _section_banner(slide, part, subtitle)
            _footer(slide, i, total)

        elif kind == "bullets":
            title, items = payload
            _title(slide, title)
            _bullets(slide, items, font_size=17)
            _footer(slide, i, total)

        elif kind == "thanks":
            tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.8),
                                          Inches(12.13), Inches(1.4))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = "Thank you"
            r.font.name = "Calibri"
            r.font.size = Pt(60)
            r.font.bold = True
            r.font.color.rgb = NAVY

            tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(4.4),
                                           Inches(12.13), Inches(0.8))
            tf2 = tb2.text_frame
            p2 = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = "Full artefacts: docs/PROJECT_OVERVIEW.{md,tex,pdf}"
            r2.font.name = "Calibri"
            r2.font.size = Pt(18)
            r2.font.color.rgb = BLUE

            tb3 = slide.shapes.add_textbox(Inches(0.6), Inches(5.2),
                                           Inches(12.13), Inches(0.6))
            tf3 = tb3.text_frame
            p3 = tf3.paragraphs[0]
            p3.alignment = PP_ALIGN.CENTER
            r3 = p3.add_run()
            r3.text = "nissimbrami@post.bgu.ac.il"
            r3.font.name = "Calibri"
            r3.font.size = Pt(14)
            r3.font.color.rgb = MUTED

    prs.save(OUT)
    print(f"Wrote {OUT} ({total} slides)")


if __name__ == "__main__":
    build()
