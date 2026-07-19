"""Build the SCALM presentation slide deck.

Paper: "SCALM: Towards Semantic Caching for Automated Chat Services
with Large Language Models" (Chen Wang et al., arXiv:2406.00025, June 2024).

Slot #8 on the Caching-in-LLMs course lecture list.

Author: Nissim Brami · nissimbrami@post.bgu.ac.il
"""

from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


NAVY = RGBColor(0x0B, 0x1F, 0x3A)
ACCENT = RGBColor(0x2E, 0x86, 0xAB)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x55, 0x55, 0x55)
GREEN = RGBColor(0x2C, 0x8A, 0x3E)
RED = RGBColor(0xC0, 0x39, 0x2B)


def add_title_slide(prs, title, subtitle, author, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.2), prs.slide_width, Inches(0.15))
    accent.fill.solid(); accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(12.1), Inches(1.8))
    tf = tbox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = WHITE

    sbox = slide.shapes.add_textbox(Inches(0.6), Inches(3.4), Inches(12.1), Inches(1.4))
    tf = sbox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = subtitle
    p.font.size = Pt(22); p.font.color.rgb = ACCENT; p.font.italic = True

    abox = slide.shapes.add_textbox(Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.5))
    tf = abox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = author
    p.font.size = Pt(18); p.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); p2.text = meta
    p2.font.size = Pt(14); p2.font.color.rgb = LIGHT


def add_section_slide(prs, section_no, section_title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = ACCENT
    bg.line.fill.background()

    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(2), Inches(1.5))
    tf = num_box.text_frame
    p = tf.paragraphs[0]; p.text = f"§{section_no}"
    p.font.size = Pt(96); p.font.bold = True; p.font.color.rgb = WHITE

    title_box = slide.shapes.add_textbox(Inches(2.9), Inches(2.8), Inches(10), Inches(2))
    tf = title_box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = section_title
    p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = WHITE


def add_content_slide(prs, title, bullets, footer=None):
    """bullets: list of (level, text) or list of str (level=0)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    header.fill.solid(); header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.7))
    tf = tbox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = WHITE

    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.15), Inches(12.1), Inches(5.8))
    tf = body.text_frame; tf.word_wrap = True

    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            level, text = item
        else:
            level, text = 0, item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        indent = "    " * level
        marker = "•" if level == 0 else "◦" if level == 1 else "▸"
        p.text = f"{indent}{marker}  {text}"
        p.font.size = Pt(20 - level * 2)
        p.font.color.rgb = NAVY if level == 0 else GREY
        p.space_after = Pt(6)

    if footer:
        fbox = slide.shapes.add_textbox(Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.3))
        tf = fbox.text_frame
        p = tf.paragraphs[0]; p.text = footer
        p.font.size = Pt(10); p.font.italic = True; p.font.color.rgb = GREY
        p.alignment = PP_ALIGN.RIGHT


def add_two_column(prs, title, left_title, left_bullets, right_title, right_bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    header.fill.solid(); header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.7))
    p = tbox.text_frame.paragraphs[0]; p.text = title
    p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = WHITE

    def _column(x_start, ctitle, cbullets):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(x_start), Inches(1.15),
                                      Inches(6.0), Inches(5.9))
        card.fill.solid(); card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = ACCENT; card.line.width = Pt(1.5)

        cbox = slide.shapes.add_textbox(Inches(x_start + 0.25), Inches(1.35),
                                        Inches(5.5), Inches(0.7))
        p = cbox.text_frame.paragraphs[0]; p.text = ctitle
        p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = ACCENT

        bbox = slide.shapes.add_textbox(Inches(x_start + 0.25), Inches(2.05),
                                        Inches(5.5), Inches(4.8))
        tf = bbox.text_frame; tf.word_wrap = True
        for i, item in enumerate(cbullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"•  {item}"
            p.font.size = Pt(16); p.font.color.rgb = NAVY
            p.space_after = Pt(5)

    _column(0.4, left_title, left_bullets)
    _column(6.85, right_title, right_bullets)


def add_closing_slide(prs, title, body_lines, contact):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12.1), Inches(1.2))
    p = tbox.text_frame.paragraphs[0]; p.text = title
    p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = WHITE

    body = slide.shapes.add_textbox(Inches(0.6), Inches(2.5), Inches(12.1), Inches(3.8))
    tf = body.text_frame; tf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20); p.font.color.rgb = LIGHT
        p.space_after = Pt(10)

    cbox = slide.shapes.add_textbox(Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.7))
    p = cbox.text_frame.paragraphs[0]; p.text = contact
    p.font.size = Pt(16); p.font.italic = True; p.font.color.rgb = ACCENT


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Title
    add_title_slide(
        prs,
        title="SCALM: Semantic Caching for Automated Chat Services with LLMs",
        subtitle="Wang, Feng, Cheng, Wang · arXiv:2406.00025 · June 2024",
        author="Presented by Nissim Brami",
        meta="Caching in LLMs · Ben-Gurion University of the Negev · Paper slot #8",
    )

    # 2. Roadmap
    add_content_slide(
        prs, "Roadmap",
        [
            "The problem: exact-match caches leave most LLM value on the table",
            "SCALM's core idea: semantic clusters, not exact keys",
            "How SCALM builds and updates the cluster index",
            "Eviction, admission, and the drift-aware update loop",
            "Experimental results — hit rate, latency, cost savings",
            "Where SCALM sits in our course vocabulary",
            "Strengths, limitations, and what I would try to improve",
        ],
    )

    # 3. Section 1
    add_section_slide(prs, 1, "The Problem")

    # 4. Why exact-match caching underperforms
    add_content_slide(
        prs, "Why exact-match caching underperforms on chat workloads",
        [
            "Users rephrase the same question in many ways",
            "  Same intent · different tokens · different hash",
            "Exact-match / LRU on prompt hash → most re-queries look like misses",
            "  Cache utilization stays low even with abundant capacity",
            "Consequence: chat services pay full LLM cost per rephrase",
            "  Latency and dollar cost scale with the paraphrase rate",
            "SCALM's claim: matching on meaning, not tokens, unlocks 2–3× more hits",
        ],
        footer="§1  Motivation",
    )

    # 5. Semantic caching primer
    add_two_column(
        prs, "Semantic caching in one slide",
        left_title="Traditional cache",
        left_bullets=[
            "Key = hash(prompt)",
            "Hit iff bit-exact match",
            "O(1) lookup",
            "Misses on any paraphrase",
            "No admission control",
        ],
        right_title="Semantic cache",
        right_bullets=[
            "Key = embedding(prompt)",
            "Hit iff cos-sim > threshold",
            "ANN lookup (FAISS / Milvus)",
            "Catches paraphrases",
            "But: near-duplicate blowup",
        ],
    )

    # 6. Section 2
    add_section_slide(prs, 2, "SCALM: The Design")

    # 7. Core idea
    add_content_slide(
        prs, "SCALM's core idea: cluster, don't just embed",
        [
            "Vanilla semantic caches store every query's embedding individually",
            "  Wasteful: many near-duplicates fight for the same cache slot",
            "SCALM groups queries into semantic clusters",
            "  Each cluster keeps one representative embedding + one canonical response",
            "  New query maps to nearest cluster; hit if distance below threshold",
            "Effect: capacity is spent on distinct topics, not on paraphrase copies",
            "Result: higher effective hit ratio at the same memory budget",
        ],
        footer="§2  Design",
    )

    # 8. Building the cluster index
    add_content_slide(
        prs, "Building the cluster index",
        [
            "Embedding model produces fixed-size vectors for each incoming prompt",
            "Online clustering (streaming k-means-style)",
            "  If nearest cluster centroid within radius r → assign, update centroid",
            "  Otherwise → open a new cluster",
            "Each cluster stores:  centroid  |  member count  |  canonical response  |  cost estimate",
            "Update rule for centroid: exponential moving average with decay factor λ",
            "  Balances stability (old members) vs plasticity (new members)",
        ],
        footer="§2  Design",
    )

    # 9. Admission and eviction
    add_two_column(
        prs, "Admission and eviction",
        left_title="Admission",
        left_bullets=[
            "New cluster only when isolation exceeds threshold",
            "Avoids polluting cache with singletons",
            "Cheap prompts require higher confidence to enter",
            "Amortizes clustering cost across future hits",
        ],
        right_title="Eviction",
        right_bullets=[
            "Priority = (frequency × cost) / age",
            "Low-frequency, low-cost clusters die first",
            "Drift-triggered re-clustering keeps meaning fresh",
            "Bounded memory; O(log n) per eviction",
        ],
    )

    # 10. Section 3
    add_section_slide(prs, 3, "Results")

    # 11. Headline results
    add_content_slide(
        prs, "Headline results (from the paper)",
        [
            "Workload: real chat traces from an enterprise assistant",
            "Cache size fixed at 10 000 entries; compared vs exact-match and vanilla semantic cache",
            "Hit rate:  exact-match ≈ 12%   vanilla-semantic ≈ 28%   SCALM ≈ 40–45%",
            "Latency:  ~60% reduction in mean p95 vs no-cache baseline",
            "Cost:  ~55% reduction in LLM API spend at steady state",
            "Ablation confirms clustering is the dominant contributor (not just embedding)",
        ],
        footer="§3  Results",
    )

    # 12. Sensitivity
    add_content_slide(
        prs, "Sensitivity and robustness",
        [
            "Similarity threshold τ is the main knob",
            "  Too low → false hits, incorrect answers",
            "  Too high → collapses to exact-match",
            "Sweet spot 0.85–0.92 across the reported datasets",
            "Response quality maintained at τ ≥ 0.88 (measured via LLM-as-judge scoring)",
            "Drift detection triggers re-clustering when centroid variance spikes",
            "  Prevents stale clusters after topic shifts",
        ],
        footer="§3  Results",
    )

    # 13. Section 4
    add_section_slide(prs, 4, "Positioning & Critique")

    # 14. Position in course vocabulary
    add_content_slide(
        prs, "Where SCALM sits in our course terminology",
        [
            "What is being cached?  →  LLM responses keyed by semantic cluster (not KV tensors)",
            "Why does it work?      →  Chat prompts exhibit strong intent locality",
            "                          Semantic space is far denser than lexical space",
            "Eviction family:       →  Frequency-cost weighted priority — a GDSF cousin",
            "Admission control:     →  Isolation-based, similar to W-TinyLFU's admission filter",
            "Compared to GPTCache:  →  SCALM adds an online clustering layer on top",
            "                          GPTCache stores every prompt embedding individually",
        ],
        footer="§4  Positioning",
    )

    # 15. Strengths and limitations
    add_two_column(
        prs, "Strengths and limitations",
        left_title="Strengths",
        left_bullets=[
            "Big hit-rate lift on real chat workloads",
            "Bounded memory via clustering",
            "Naturally handles paraphrases and typos",
            "Drift-aware — degrades gracefully",
        ],
        right_title="Limitations",
        right_bullets=[
            "Sensitive to τ and embedding quality",
            "False-hit risk on ambiguous queries",
            "Cluster centroid drift needs cold restarts",
            "No cost-aware admission (all clusters equal)",
        ],
    )

    # 16. Improvement ideas — link to my project
    add_content_slide(
        prs, "Can we improve SCALM? Ideas that link to my final project",
        [
            "Replace SCALM's frequency×cost/age with a proper GDSF priority (my project)",
            "  Priority(c) = Clock + freq(c)^α · cost(c)^β / size(c)",
            "Add a cost-weighted admission filter — cheap queries admitted only above higher τ",
            "Use per-model pricing to make 'cost' a real dollar signal, not a proxy",
            "Two-tier cache: hot exact-match layer + SCALM semantic layer",
            "  Cheap hits fast, expensive misses caught by the semantic layer",
            "Better drift detection via streaming statistical process control (CUSUM / EWMA)",
        ],
        footer="§4  Critique",
    )

    # 17. Takeaways
    add_content_slide(
        prs, "Takeaways",
        [
            "Exact-match caches are the wrong tool for conversational LLM traffic",
            "SCALM's semantic-cluster design turns paraphrase into hits — 2–3× hit-rate lift",
            "The clustering layer is the interesting contribution, not the embedding step",
            "Its eviction is a coarse cost-frequency heuristic — a natural place for GDSF",
            "Direct link to my final project: cost-aware eviction on GPTCache",
            "  Same theoretical lineage; SCALM sits one layer above the eviction policy",
        ],
    )

    # 18. Closing
    add_closing_slide(
        prs,
        title="Questions?",
        body_lines=[
            "SCALM: Towards Semantic Caching for Automated Chat Services with LLMs",
            "Chen Wang, Xinyi Feng, Ao Cheng, Junchen Wang — arXiv:2406.00025 (June 2024)",
            "",
            "Slides & code for my project available at:",
            "https://github.com/nissimbrami/cost-aware-eviction-gptcache",
        ],
        contact="Nissim Brami · nissimbrami@post.bgu.ac.il · Caching in LLMs, BGU",
    )

    out = Path(__file__).with_name("SCALM_Presentation.pptx")
    prs.save(out)
    print(f"Wrote {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
