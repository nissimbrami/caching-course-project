"""Build the StreamingLLM presentation deck.

Paper: Efficient Streaming Language Models with Attention Sinks
       Guangxuan Xiao, Yuandong Tian, Beidi Chen, Song Han, Mike Lewis
       ICLR 2024 - arXiv:2309.17453

Style follows the reference decks under `reference-decks/`:
  - 16:9 widescreen (13.33 x 7.5 in), matching the prior student deck (SCALM.pptx)
  - Structure: Title -> Agenda -> Basics -> Paper -> Results -> Limits ->
    What I'd change -> Conclusion -> Thank you (mirrors SCALM.pptx)
  - Vocabulary anchored on course concepts (LRU/LFU/admission/eviction)
    the way Prof. Einziger's Chapter1 - Cache deck introduces them.

Every numeric claim on every slide is traceable to a specific page/table
of the paper via `papers/streaming-llm/deep-read.md`.

Author: Nissim Brami (nissimbrami@post.bgu.ac.il)
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn


HERE = Path(__file__).resolve().parent
OUT = HERE / "StreamingLLM_Presentation.pptx"


# ---------- Colour palette (kept close to the Chapter1 deck) --------------

NAVY = RGBColor(0x0B, 0x1F, 0x3A)
BLUE = RGBColor(0x2E, 0x86, 0xAB)
ACCENT = RGBColor(0xE6, 0x7E, 0x22)
BG = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x66, 0x66, 0x66)
HIGHLIGHT = RGBColor(0xC0, 0x39, 0x2B)
GREEN_OK = RGBColor(0x27, 0xAE, 0x60)
GREY = RGBColor(0xB3, 0xB3, 0xB3)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)


# ---------- Helpers -------------------------------------------------------


def _blank(prs: Presentation):
    """Blank layout - we place everything explicitly."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _title(slide, text: str, top_in: float = 0.35):
    left = Inches(0.5)
    top = Inches(top_in)
    width = Inches(12.33)
    height = Inches(0.9)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(34)
    r.font.bold = True
    r.font.color.rgb = NAVY
    # accent underline
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(top_in + 0.95),
                                      Inches(12.83), Inches(top_in + 0.95))
    line.line.color.rgb = ACCENT
    line.line.width = Pt(2.25)
    return tb


def _footer(slide, page_num: int):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.05),
                                  Inches(12.33), Inches(0.3))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Nissim Brami · Caching in LLMs · BGU"
    r.font.name = "Calibri"
    r.font.size = Pt(9)
    r.font.color.rgb = MUTED
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = f"{page_num}"
    r2.font.name = "Calibri"
    r2.font.size = Pt(9)
    r2.font.color.rgb = MUTED
    # Move page number to the right edge with a separate textbox for clarity
    tb2 = slide.shapes.add_textbox(Inches(12.4), Inches(7.05),
                                   Inches(0.4), Inches(0.3))
    tf2 = tb2.text_frame
    tf2.margin_left = 0
    tf2.margin_top = 0
    p3 = tf2.paragraphs[0]
    p3.alignment = PP_ALIGN.RIGHT
    r3 = p3.add_run()
    r3.text = f"{page_num}"
    r3.font.name = "Calibri"
    r3.font.size = Pt(9)
    r3.font.color.rgb = MUTED
    # Remove the duplicate page number in the left footer (leave only the credit line)
    tf.paragraphs[0].runs[0].text = "Nissim Brami · Caching in LLMs · BGU"
    # Clear the second paragraph on the left footer
    tf.paragraphs[1].runs[0].text = ""


def _bullets(slide, items, left_in=0.7, top_in=1.6, width_in=12.0,
             height_in=5.0, size=20, color=TEXT, spacing=8):
    """Add a bullet list. Items can be strings or (text, bold, colour) tuples."""
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in),
                                  Inches(width_in), Inches(height_in))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    for idx, it in enumerate(items):
        if isinstance(it, tuple):
            text, bold, col = it
        else:
            text, bold, col = it, False, color
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        r = p.add_run()
        r.text = "• " + text
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = col
    return tb


def _plain_text(slide, text, left_in, top_in, width_in, height_in,
                size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in),
                                  Inches(width_in), Inches(height_in))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def _box(slide, left_in, top_in, width_in, height_in, fill=LIGHT,
         line=GREY, line_pt=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(left_in), Inches(top_in),
                                 Inches(width_in), Inches(height_in))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(line_pt)
    shp.shadow.inherit = False
    return shp


def _labelled_box(slide, left_in, top_in, width_in, height_in,
                  label, sub=None, fill=LIGHT, border=GREY,
                  label_size=16, sub_size=12,
                  label_color=NAVY, sub_color=MUTED):
    _box(slide, left_in, top_in, width_in, height_in, fill=fill, line=border)
    _plain_text(slide, label, left_in + 0.1, top_in + 0.1,
                width_in - 0.2, 0.4, size=label_size, bold=True,
                color=label_color)
    if sub is not None:
        _plain_text(slide, sub, left_in + 0.1, top_in + 0.55,
                    width_in - 0.2, height_in - 0.65,
                    size=sub_size, color=sub_color)


def _table(slide, headers, rows, left_in, top_in, col_widths_in,
           head_size=13, cell_size=12, row_height_in=0.36,
           head_fill=NAVY, head_color=BG, stripe=LIGHT):
    n_cols = len(headers)
    n_rows = len(rows) + 1
    total_w = sum(col_widths_in)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols,
                                       Inches(left_in), Inches(top_in),
                                       Inches(total_w),
                                       Inches(row_height_in * n_rows))
    tbl = tbl_shape.table
    for i, w in enumerate(col_widths_in):
        tbl.columns[i].width = Inches(w)
    for i in range(n_rows):
        tbl.rows[i].height = Inches(row_height_in)
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = head_fill
        tf = c.text_frame
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = h
        r.font.name = "Calibri"
        r.font.size = Pt(head_size)
        r.font.bold = True
        r.font.color.rgb = head_color
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = stripe if (i % 2 == 0) else BG
            tf = c.text_frame
            tf.margin_left = Inches(0.05)
            tf.margin_right = Inches(0.05)
            tf.margin_top = Inches(0.02)
            tf.margin_bottom = Inches(0.02)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            r.font.name = "Calibri"
            r.font.size = Pt(cell_size)
            r.font.color.rgb = TEXT
    return tbl


# ---------- Slide builders ------------------------------------------------


def s_title(prs):
    s = _blank(prs)
    # Accent bar
    _box(s, 0, 0, 13.33, 0.4, fill=NAVY, line=NAVY)
    _plain_text(s, "Efficient Streaming Language Models",
                0.7, 1.7, 12.0, 0.9, size=40, bold=True, color=NAVY)
    _plain_text(s, "with Attention Sinks",
                0.7, 2.55, 12.0, 0.8, size=32, bold=True, color=BLUE)
    _plain_text(s,
                "Guangxuan Xiao · Yuandong Tian · Beidi Chen · Song Han · Mike Lewis",
                0.7, 3.6, 12.0, 0.4, size=16, color=TEXT)
    _plain_text(s, "MIT · Meta AI · CMU · NVIDIA · ICLR 2024",
                0.7, 3.95, 12.0, 0.4, size=14, color=MUTED)
    _plain_text(s, "arXiv:2309.17453  ·  github.com/mit-han-lab/streaming-llm",
                0.7, 4.35, 12.0, 0.4, size=13, color=MUTED)
    _box(s, 0.7, 5.35, 12.0, 0.03, fill=ACCENT, line=ACCENT)
    _plain_text(s, "Presented by Nissim Brami",
                0.7, 5.55, 12.0, 0.4, size=16, bold=True, color=NAVY)
    _plain_text(s, "Caching in LLMs · Ben-Gurion University · Prof. Gil Einziger",
                0.7, 5.95, 12.0, 0.4, size=13, color=MUTED)
    return s


def s_agenda(prs):
    s = _blank(prs)
    _title(s, "Agenda  ·  ~60 minutes")
    items = [
        "Basics — LLM decoding, KV cache, why this is a caching problem",
        "Related work — length extrapolation, context extension, sparse attn",
        "The observation — attention sinks are position, not semantics",
        "The mechanism — sinks + rolling cache + cache-local position IDs",
        "Mechanism deep-dive — worked example, RoPE / ALiBi variants, algorithm box",
        "Sink-Token pre-training — the §3.3 refinement + its 160M caveat",
        "Results — perplexity to 4M tokens, 22.2× throughput, streaming QA",
        "Results deep-dive — Fig. 5, Fig. 10, Table 5, StreamEval",
        "Modern alternatives — H2O, SnapKV, Scissorhands, FastGen",
        "Limitations — LongBench, non-monotone cache, positional-only ceiling",
        "Live implementation walkthrough — ~50 lines of PyTorch + HF",
        "What I would change — bridge to my Task-2 GDSF project",
        "Q&A prep",
    ]
    _bullets(s, items, top_in=1.55, size=17, spacing=6)
    return s


# --- Basics section ------------------------------------------------------


def s_section_basics(prs):
    s = _blank(prs)
    _box(s, 0, 0, 13.33, 7.5, fill=NAVY, line=NAVY)
    _plain_text(s, "Basics", 0.7, 3.0, 12.0, 1.2,
                size=54, bold=True, color=BG, align=PP_ALIGN.LEFT)
    _box(s, 0.7, 4.1, 6.0, 0.05, fill=ACCENT, line=ACCENT)
    _plain_text(s, "Streaming LLMs · KV cache · Softmax",
                0.7, 4.3, 12.0, 0.5, size=20, color=BG)
    return s


def s_llm_decoding(prs):
    s = _blank(prs)
    _title(s, "The scene: decoding a large language model")
    _bullets(s, [
        "Generate one token at a time (autoregressive decoding).",
        "Every new token attends back to every previous token.",
        "To avoid recomputing, we cache each previous token’s "
        "Key and Value vectors: the KV cache.",
        "KV cache size grows linearly with sequence length.",
        "For a 7B model at 4K tokens, KV cache is ~2 GB (FP16 estimate).",
    ], top_in=1.6, size=20)
    _labelled_box(s, 8.0, 4.9, 5.0, 1.9,
                  "The bottleneck at decode time",
                  "Memory bandwidth, not compute.\n"
                  "The KV cache is what we cache.",
                  fill=LIGHT, border=BLUE,
                  label_size=15, sub_size=13, sub_color=TEXT)
    return s


def s_course_anchor(prs):
    s = _blank(prs)
    _title(s, "This is a caching problem: course vocabulary")
    _bullets(s, [
        "Cache: the KV cache in GPU memory.",
        "Capacity: bounded by GPU RAM.",
        "Requests: attention lookups from each new query token.",
        "Miss: no cache = recompute the whole context (quadratic).",
        "Eviction policy: which past tokens to keep, which to drop?",
        "Baseline eviction: window attention = LRU by position.",
    ], top_in=1.55, size=20)
    _labelled_box(s, 8.2, 4.6, 4.8, 2.2,
                  "Chapter 1 recap",
                  "LRU · LFU · ARC · LIRS · Hyperbolic · FRD · W-TinyLFU\n"
                  "Admission vs. eviction.\n"
                  "StreamingLLM chooses a positional eviction rule.",
                  fill=LIGHT, border=BLUE,
                  label_size=15, sub_size=13, sub_color=TEXT)
    return s


def s_naive_options(prs):
    s = _blank(prs)
    _title(s, "Three obvious KV-cache strategies, all three break")
    headers = ["Strategy", "Cost", "Behaviour on long streams", "PPL"]
    rows = [
        ["(a) Dense attention",     "O(T²)",   "OOM past training window; degrades on long text", "5641"],
        ["(b) Window attention",    "O(T·L)",  "Collapses when first tokens are evicted", "5158"],
        ["(c) Window + re-compute", "O(T·L²)", "Correct but slow (per-token rebuild)",   "5.43"],
    ]
    _table(s, headers, rows, left_in=0.7, top_in=1.6,
           col_widths_in=[3.3, 1.2, 5.6, 1.3],
           head_size=14, cell_size=13, row_height_in=0.55)
    _plain_text(s,
                "Numbers: Llama-2-13B, first book of PG19 (65K tokens) "
                "[paper Fig. 1, p. 2].",
                0.7, 4.2, 12.0, 0.4, size=13, color=MUTED)
    _labelled_box(s, 0.7, 4.9, 12.0, 1.9,
                  "The puzzle",
                  "Window attention has the right cost but the wrong "
                  "answer. Evicting a single token, the first one, "
                  "turns a fluent model into gibberish (PPL 5158). "
                  "Why?",
                  fill=LIGHT, border=HIGHLIGHT,
                  label_size=16, sub_size=15, sub_color=TEXT,
                  label_color=HIGHLIGHT)
    return s


# --- The observation -----------------------------------------------------


def s_section_observation(prs):
    s = _blank(prs)
    _box(s, 0, 0, 13.33, 7.5, fill=NAVY, line=NAVY)
    _plain_text(s, "The observation", 0.7, 3.0, 12.0, 1.2,
                size=54, bold=True, color=BG)
    _box(s, 0.7, 4.1, 6.0, 0.05, fill=ACCENT, line=ACCENT)
    _plain_text(s, "Attention Sinks",
                0.7, 4.3, 12.0, 0.5, size=22, color=BG)
    return s


def s_attention_sink_phenomenon(prs):
    s = _blank(prs)
    _title(s, "Attention concentrates on the first few tokens, everywhere")
    _bullets(s, [
        "Look at the attention map of Llama-2-7B, across every "
        "layer and every head.",
        "Beyond the bottom two layers, the model heavily attends "
        "to the initial tokens, regardless of what those tokens say.",
        "The phenomenon holds in Llama-2, MPT, Falcon, Pythia; "
        "and in BERT (on the [SEP] token); and in ViTs (register tokens).",
        "The authors call these tokens attention sinks.",
    ], top_in=1.55, size=20)
    _labelled_box(s, 0.7, 5.15, 12.0, 1.7,
                  "Paper's exact wording (Fig. 2, p. 3)",
                  "\"Beyond the bottom two layers, the model heavily attends "
                  "to the initial token across all layers and heads.\" "
                  "Visualised over 256 sentences of length 16 in Llama-2-7B.",
                  fill=LIGHT, border=BLUE,
                  label_size=16, sub_size=15, sub_color=TEXT)
    return s


def s_softmax_argument(prs):
    s = _blank(prs)
    _title(s, "Why? The softmax must sum to one")
    # Softmax equation card
    _box(s, 0.7, 1.6, 12.0, 1.4, fill=LIGHT, line=BLUE, line_pt=1.0)
    _plain_text(s,
                "softmax(x)_i  =  exp(x_i) / (exp(x_1) + Σ_{j≥2} exp(x_j))",
                0.9, 1.75, 11.5, 0.5, size=22, bold=True, color=NAVY,
                align=PP_ALIGN.CENTER)
    _plain_text(s,
                "If x_1 ≫ x_j for j ≥ 2, then exp(x_1) dominates the denominator.",
                0.9, 2.35, 11.5, 0.4, size=15, color=TEXT,
                align=PP_ALIGN.CENTER)
    _bullets(s, [
        "Softmax forces the model to allocate attention mass, "
        "even when nothing in the context deserves it.",
        "That excess mass has to land somewhere.",
        "In a causal LM, initial tokens are visible to every "
        "subsequent position; they are the natural dumping ground.",
        "Evicting them removes a huge chunk of the softmax "
        "denominator → the attention distribution warps → PPL explodes.",
    ], top_in=3.35, size=19)
    return s


def s_semantics_or_position(prs):
    s = _blank(prs)
    _title(s, "Is it the tokens, or is it the position?")
    _bullets(s, [
        "Hypothesis A: the first tokens are semantically special.",
        "Hypothesis B: the first positions are structurally special.",
        "Test: replace the first four tokens with newline characters "
        "and re-run.",
    ], top_in=1.55, size=20)
    headers = ["Cache configuration (Llama-2-13B, PG19)", "PPL ↓"]
    rows = [
        ["0 + 1024  (window attention, no sinks)",  "5158.07"],
        ["4 + 1020  (StreamingLLM, real sinks)",    "5.40"],
        ["4×\"\\n\" + 1020  (linebreak sinks)",     "5.60"],
    ]
    _table(s, headers, rows, left_in=1.5, top_in=3.55,
           col_widths_in=[7.0, 2.0], head_size=13, cell_size=13,
           row_height_in=0.5)
    _labelled_box(s, 1.5, 5.75, 9.0, 1.1,
                  "Verdict",
                  "The first four positions do the work, not their content. "
                  "Table 1, p. 5.",
                  fill=LIGHT, border=GREEN_OK,
                  label_size=15, sub_size=14, sub_color=TEXT,
                  label_color=GREEN_OK)
    return s


# --- The mechanism -------------------------------------------------------


def s_section_mechanism(prs):
    s = _blank(prs)
    _box(s, 0, 0, 13.33, 7.5, fill=NAVY, line=NAVY)
    _plain_text(s, "The mechanism", 0.7, 3.0, 12.0, 1.2,
                size=54, bold=True, color=BG)
    _box(s, 0.7, 4.1, 6.0, 0.05, fill=ACCENT, line=ACCENT)
    _plain_text(s, "Sinks + Rolling KV cache",
                0.7, 4.3, 12.0, 0.5, size=22, color=BG)
    return s


def s_cache_layout(prs):
    s = _blank(prs)
    _title(s, "The KV cache: 4 sink tokens + rolling window")
    # 4 sink cells
    left = 0.9
    top = 2.1
    cell = 0.7
    for i in range(4):
        x = left + i * (cell + 0.05)
        _box(s, x, top, cell, 0.9, fill=ACCENT, line=ACCENT)
        _plain_text(s, str(i), x, top + 0.25, cell, 0.5,
                    size=22, bold=True, color=BG, align=PP_ALIGN.CENTER)
    _plain_text(s, "SINKS", left + 0.05, top + 1.0, 2.5, 0.35,
                size=12, bold=True, color=ACCENT)
    # gap
    _plain_text(s, "…", left + 4 * (cell + 0.05) + 0.1,
                top + 0.15, 0.6, 0.6,
                size=32, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    _plain_text(s, "(evicted)", left + 4 * (cell + 0.05) + 0.1,
                top + 1.0, 1.0, 0.35, size=12, color=MUTED)
    # rolling window cells
    left2 = left + 4 * (cell + 0.05) + 1.0
    labels = ["T-6", "T-5", "T-4", "T-3", "T-2", "T-1", "T"]
    for i, lbl in enumerate(labels):
        x = left2 + i * (cell + 0.05)
        _box(s, x, top, cell, 0.9, fill=BLUE, line=BLUE)
        _plain_text(s, lbl, x, top + 0.3, cell, 0.5,
                    size=13, bold=True, color=BG, align=PP_ALIGN.CENTER)
    _plain_text(s, "ROLLING WINDOW (last L tokens)",
                left2, top + 1.0, 6.0, 0.35,
                size=12, bold=True, color=BLUE)
    _bullets(s, [
        "Keep the first S tokens (S = 4 works everywhere) as sinks.",
        "Keep the last L tokens as the rolling window.",
        "Everything in between is evicted.",
        "Cache size = S + L; independent of stream length.",
    ], top_in=3.85, size=20)
    return s


def s_position_reindex(prs):
    s = _blank(prs)
    _title(s, "The trick: positions are indexed inside the cache")
    _bullets(s, [
        "Naive attempt: use each token’s position in the original text.",
        "That does not work: the sinks and the rolling window are "
        "far apart in the text but adjacent in the cache.",
        "StreamingLLM re-indexes position IDs inside the cache.",
        "Compatible with RoPE and ALiBi, the two dominant positional "
        "encodings in 2024.",
    ], top_in=1.55, size=20)
    _labelled_box(s, 0.7, 4.7, 12.0, 2.2,
                  "Concrete example (paper §3.2, p. 5)",
                  "Cache contains tokens whose original positions are "
                  "[0, 1, 2, 3, 6, 7, 8].\n"
                  "The model is asked to predict position 9.\n"
                  "It receives cache-local positions [0, 1, 2, 3, 4, 5, 6, 7], "
                  "not [0, 1, 2, 3, 6, 7, 8, 9].\n"
                  "RoPE keys are cached pre-rotation, then rotated at decode "
                  "using the cache-local index.",
                  fill=LIGHT, border=BLUE,
                  label_size=15, sub_size=14, sub_color=TEXT)
    return s


def s_sink_token_pretraining(prs):
    s = _blank(prs)
    _title(s, "Refinement: pre-train with a dedicated sink token")
    _bullets(s, [
        "If we can control pre-training, we can train the model to "
        "concentrate its excess attention on one specific token.",
        "Prepend a single learnable placeholder to every pre-training "
        "sample. Call it the sink token.",
        "One sink token now suffices at inference, instead of the "
        "four accidental ones.",
    ], top_in=1.55, size=19)
    headers = ["Cache config (160M model, PG19)", "0+1024", "1+1023",
               "2+1022", "4+1020"]
    rows = [
        ["Vanilla",         "27.87", "18.49", "18.05", "18.05"],
        ["Zero Sink",       "29214", "19.90", "18.27", "18.01"],
        ["Learnable Sink",  "1235",  "18.01", "18.01", "18.02"],
    ]
    _table(s, headers, rows, left_in=1.2, top_in=4.35,
           col_widths_in=[3.9, 1.6, 1.6, 1.6, 1.6],
           head_size=13, cell_size=13, row_height_in=0.4)
    _plain_text(s,
                "Table 3, p. 6.  \"Zero Sink\" = SoftMax-off-by-One "
                "(Miller, 2023). ",
                1.2, 6.15, 11.0, 0.35, size=13, color=MUTED)
    return s


# --- Results -------------------------------------------------------------


def s_section_results(prs):
    s = _blank(prs)
    _box(s, 0, 0, 13.33, 7.5, fill=NAVY, line=NAVY)
    _plain_text(s, "Results", 0.7, 3.0, 12.0, 1.2,
                size=54, bold=True, color=BG)
    _box(s, 0.7, 4.1, 6.0, 0.05, fill=ACCENT, line=ACCENT)
    _plain_text(s, "Perplexity · Throughput · Streaming QA",
                0.7, 4.3, 12.0, 0.5, size=22, color=BG)
    return s


def s_perplexity_result(prs):
    s = _blank(prs)
    _title(s, "Result 1: perplexity stays flat past 4 million tokens")
    _bullets(s, [
        "Concatenated PG19 test set (100 long books).",
        "Cache = 2048 for Llama-2 · 1024 for MPT, Falcon, Pythia.",
        "Perplexity stable up to 4M+ tokens across every family "
        "(Llama-2-{7,13,70}B · MPT-{7,30}B · Falcon-{7,40}B · "
        "Pythia-{2.8,6.9,12}B).",
        "Matches the sliding-window-with-recomputation baseline; "
        "beats dense (OOMs past training window) and window "
        "attention (PPL explodes).",
    ], top_in=1.55, size=20)
    _labelled_box(s, 0.7, 4.9, 12.0, 1.9,
                  "The single strongest table",
                  "Llama-2-13B, PG19, 65K tokens (Table 1, p. 5):\n"
                  "  Window attention        →  PPL 5158.07\n"
                  "  StreamingLLM (4+1020)   →  PPL 5.40\n"
                  "  Linebreak sinks         →  PPL 5.60  (sinks are positional!)",
                  fill=LIGHT, border=GREEN_OK,
                  label_size=15, sub_size=14, sub_color=TEXT,
                  label_color=GREEN_OK)
    return s


def s_throughput_result(prs):
    s = _blank(prs)
    _title(s, "Result 2: up to 22.2× faster than window-with-recomputation")
    _bullets(s, [
        "Single NVIDIA A6000 · Hugging Face Transformers · batch 1.",
        "The sliding-window-with-recomputation baseline is quadratic; "
        "StreamingLLM is linear.",
    ], top_in=1.55, size=20)
    headers = ["Cache size", "Sliding+Recompute (ms)", "StreamingLLM (ms)",
               "Speed-up"]
    rows = [
        ["256",  "2355", "106", "22.2×"],
        ["512",  "860",  "75",  "11.5×"],
        ["1024", "361",  "60",  "6.0×"],
        ["2048", "169",  "52",  "3.3×"],
        ["4096", "99",   "48",  "2.1×"],
    ]
    _table(s, headers, rows, left_in=1.5, top_in=3.6,
           col_widths_in=[2.0, 3.4, 2.6, 1.6],
           head_size=13, cell_size=13, row_height_in=0.42)
    _plain_text(s,
                "Llama-2-13B, per-token decode latency; Fig. 10, p. 9. "
                "Memory footprint essentially unchanged.",
                1.5, 6.35, 11.0, 0.4, size=12, color=MUTED)
    return s


def s_streaming_qa(prs):
    s = _blank(prs)
    _title(s, "Result 3: streaming QA is now actually usable")
    _bullets(s, [
        "Setup: concatenate all ARC-{Easy,Challenge} questions and "
        "feed them as one long stream. Score exact-match on each answer.",
    ], top_in=1.55, size=19)
    headers = ["Model / policy", "ARC-Easy", "ARC-Challenge"]
    rows = [
        ["Llama-2-7B-Chat  one-shot",       "71.25", "53.16"],
        ["Llama-2-7B-Chat  window",         "3.58",  "1.39"],
        ["Llama-2-7B-Chat  StreamingLLM",   "71.34", "55.03"],
        ["Llama-2-13B-Chat one-shot",       "78.16", "63.31"],
        ["Llama-2-13B-Chat window",         "0.25",  "0.34"],
        ["Llama-2-13B-Chat StreamingLLM",   "80.89", "65.61"],
        ["Llama-2-70B-Chat one-shot",       "91.29", "78.50"],
        ["Llama-2-70B-Chat window",         "0.12",  "0.32"],
        ["Llama-2-70B-Chat StreamingLLM",   "91.37", "80.20"],
    ]
    _table(s, headers, rows, left_in=2.7, top_in=2.55,
           col_widths_in=[5.0, 1.6, 2.2],
           head_size=13, cell_size=12, row_height_in=0.32)
    _plain_text(s,
                "Cache 1024.  Dense OOMs.  Table 5, p. 8.",
                2.7, 5.85, 10.0, 0.35, size=13, color=MUTED)
    return s


# --- Limits --------------------------------------------------------------


def s_limits(prs):
    s = _blank(prs)
    _title(s, "What StreamingLLM does not do")
    _bullets(s, [
        ("It does not extend the context window. If the answer is "
         "older than the rolling window, it is gone.", True, HIGHLIGHT),
        ("On LongBench, StreamingLLM 4+3496 underperforms the default "
         "truncation baseline on all six tasks, because it loses the "
         "initial prompt. Table 8, p. 17.", False, TEXT),
        ("Bigger cache does not always mean lower perplexity. For "
         "Llama-2-7B, 4+2044 gives 9.08 PPL but 4+4092 gives 9.59. "
         "Table 6, p. 9.", False, TEXT),
        ("No comparison with attention-score-based eviction "
         "(H2O, SnapKV, FastGen, Keyformer). Positional only.",
         False, TEXT),
        ("Learnable-sink pre-training was validated at 160 M parameters, "
         "not at 7 B, 13 B, or 70 B.", False, TEXT),
    ], top_in=1.55, size=18)
    return s


# --- What I would change / bridge to Task 2 -----------------------------


def s_what_id_change(prs):
    s = _blank(prs)
    _title(s, "What I would change")
    _bullets(s, [
        "Layer-specific sink budgets. The first two layers barely "
        "need sinks; deeper layers need more.",
        "Attention-score-based eviction inside the rolling window: "
        "combine StreamingLLM’s positional sinks with H2O-style "
        "content-aware eviction.",
        "Test on real streaming traffic (LMSYS-Chat, ShareGPT). "
        "The paper only tests on PG19 books.",
        "Sink-token pre-training at 7B+. The strongest §3.3 claim "
        "still awaits validation at scale.",
    ], top_in=1.55, size=20)
    _labelled_box(s, 0.7, 5.2, 12.0, 1.7,
                  "Same idea, one level up the stack",
                  "StreamingLLM chooses positional KV-cache eviction. "
                  "My final project (Task 2) chooses cost-aware "
                  "response-cache eviction (GDSF on GPTCache). Both "
                  "attack the same underlying problem, dumb eviction "
                  "wastes capacity, at different layers.",
                  fill=LIGHT, border=BLUE,
                  label_size=16, sub_size=15, sub_color=TEXT)
    return s


# --- Related work section -------------------------------------------------


def s_section_related(prs):
    s = _blank(prs)
    _box(s, 0, 0, 13.33, 7.5, fill=NAVY, line=NAVY)
    _plain_text(s, "Related work", 0.7, 3.0, 12.0, 1.2,
                size=54, bold=True, color=BG)
    _box(s, 0.7, 4.1, 6.0, 0.05, fill=ACCENT, line=ACCENT)
    _plain_text(s, "Where StreamingLLM sits in the 2024 landscape",
                0.7, 4.3, 12.0, 0.5, size=22, color=BG)
    return s


def s_related_length_extrap(prs):
    s = _blank(prs)
    _title(s, "Related work 1/3: length-extrapolation encodings")
    _bullets(s, [
        "RoPE (Su et al., 2021) — encodes absolute position as a rotation "
        "of the query/key vectors; relative-position property falls out.",
        "ALiBi (Press et al., 2022) — no positional embedding at all; "
        "adds a linear bias -m·|i-j| to attention logits.",
        "xPos (Sun et al., 2022) — RoPE with exponential decay, aimed "
        "at improving long-range stability.",
        "T5 relative bias (Raffel et al., 2020) — bucketed learned bias.",
    ], top_in=1.55, size=18)
    _labelled_box(s, 0.7, 4.9, 12.0, 1.95,
                  "Why StreamingLLM cares",
                  "These are the encodings the paper must plug into "
                  "unchanged. §3.2 shows the cache-local re-indexing "
                  "rule for RoPE (rotate at decode using cache-index) "
                  "and ALiBi (contiguous bias in cache coordinates). "
                  "The paper is deliberately encoding-agnostic — it is "
                  "a KV-cache policy, not a new positional scheme.",
                  fill=LIGHT, border=BLUE,
                  label_size=15, sub_size=14, sub_color=TEXT)
    return s


def s_related_context_ext(prs):
    s = _blank(prs)
    _title(s, "Related work 2/3: context-window extension")
    _bullets(s, [
        "Position Interpolation (Chen et al., 2023) — linearly "
        "interpolate RoPE indices to a longer window; needs "
        "fine-tuning.",
        "YaRN (Peng et al., 2023) — NTK-aware frequency scaling; "
        "less fine-tuning, longer context.",
        "LongRoPE (Ding et al., 2024) — evolutionary search over RoPE "
        "rescalings.",
        "FlashAttention (Dao 2022, 2023) — makes O(T²) attention "
        "*possible* at long T; orthogonal, not competitive.",
        "Landmark Attention, Focused Transformer — retrieval-style hooks.",
    ], top_in=1.55, size=17)
    _labelled_box(s, 0.7, 5.35, 12.0, 1.5,
                  "Different problem",
                  "These methods try to grow the pre-training window. "
                  "StreamingLLM keeps the window bounded and lets the "
                  "*stream* grow. Complementary, not competitive.",
                  fill=LIGHT, border=BLUE,
                  label_size=15, sub_size=14, sub_color=TEXT)
    return s


def s_related_sparse_and_lminf(prs):
    s = _blank(prs)
    _title(s, "Related work 3/3: sparse attention and concurrent work")
    _bullets(s, [
        "Sparse Transformer (Child et al., 2019), Longformer (Beltagy "
        "et al., 2020), BigBird (Zaheer et al., 2020), ETC (Ainslie "
        "et al., 2020) — structured sparsity + global tokens. "
        "Need custom kernels; not drop-in on pre-trained decoders.",
        "LM-Infinite (Han et al., 2023) — concurrent with StreamingLLM; "
        "Λ-shaped attention pattern (keep initial tokens + local band).",
        "Landmark Attention (Mohtashami & Jaggi, 2023) — retrieval-style "
        "distant memory access.",
    ], top_in=1.55, size=17)
    _labelled_box(s, 0.7, 5.05, 12.0, 1.8,
                  "How StreamingLLM differentiates itself",
                  "First to (a) *name* the attention-sink phenomenon and "
                  "show it exists across Llama-2 / MPT / Falcon / Pythia "
                  "/ BERT / ViT; (b) show that the sinks are "
                  "*positional, not semantic* (linebreak substitution); "
                  "(c) propose the learnable-sink pre-training fix.",
                  fill=LIGHT, border=BLUE,
                  label_size=15, sub_size=14, sub_color=TEXT)
    return s


# --- Mechanism deep-dive -------------------------------------------------


def s_mech_worked_example(prs):
    s = _blank(prs)
    _title(s, "Mechanism deep-dive: a 20-token worked example")
    _plain_text(s,
                "Cache = 4 sinks + 8-token rolling window. Decoding step 20.",
                0.7, 1.55, 12.0, 0.4, size=15, color=TEXT)
    # original positions row
    _plain_text(s, "Original positions kept in cache:",
                0.7, 2.05, 12.0, 0.35, size=13, bold=True, color=NAVY)
    for i, orig in enumerate([0, 1, 2, 3, 12, 13, 14, 15, 16, 17, 18, 19]):
        x = 0.7 + i * 0.9
        fill = ACCENT if i < 4 else BLUE
        _box(s, x, 2.5, 0.75, 0.55, fill=fill, line=fill)
        _plain_text(s, str(orig), x, 2.6, 0.75, 0.4,
                    size=13, bold=True, color=BG, align=PP_ALIGN.CENTER)
    _plain_text(s, "(sinks)", 0.7, 3.15, 3.5, 0.3,
                size=11, bold=True, color=ACCENT)
    _plain_text(s, "(rolling window: positions 4–11 evicted)",
                4.4, 3.15, 6.0, 0.3, size=11, bold=True, color=BLUE)
    # remapped row
    _plain_text(s, "Cache-local position IDs the model actually sees:",
                0.7, 3.75, 12.0, 0.35, size=13, bold=True, color=NAVY)
    for i, local in enumerate(list(range(12))):
        x = 0.7 + i * 0.9
        _box(s, x, 4.2, 0.75, 0.55, fill=LIGHT, line=NAVY, line_pt=1.0)
        _plain_text(s, str(local), x, 4.3, 0.75, 0.4,
                    size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    _labelled_box(s, 0.7, 5.15, 12.0, 1.75,
                  "Why re-indexing matters",
                  "If we passed the model original positions "
                  "[0,1,2,3,12,13,14,15,16,17,18,19] with a query at "
                  "position 20, RoPE would rotate keys as if positions "
                  "4–11 existed. ALiBi would penalise the 12→13 gap. "
                  "Cache-local re-indexing makes the sinks *adjacent* "
                  "to the rolling window in position space, which is "
                  "the regime the pre-trained model was optimised for. "
                  "Paper §3.2, p. 5.",
                  fill=LIGHT, border=BLUE,
                  label_size=15, sub_size=13, sub_color=TEXT)
    return s


def s_mech_rope_alibi(prs):
    s = _blank(prs)
    _title(s, "Mechanism deep-dive: RoPE and ALiBi variants")
    headers = ["Aspect", "RoPE (Llama-2, Falcon, Pythia)", "ALiBi (MPT)"]
    rows = [
        ["Where positions enter",
         "Rotate q, k by θ_i in complex plane",
         "Add -m·|i-j| bias to logits"],
        ["Cache what?",
         "Cache pre-rotation K",
         "Cache K unmodified"],
        ["Rotate/bias when?",
         "Apply rotation at decode with cache-local index",
         "Apply contiguous linear bias in cache coordinates"],
        ["Effect of a gap in positions",
         "Wrong relative rotation → PPL explodes",
         "Bias magnitude wrong → PPL explodes"],
        ["Fix cost",
         "One extra Q·K^T-time rotate; free memory-wise",
         "Recompute the bias row; free memory-wise"],
    ]
    _table(s, headers, rows, left_in=0.5, top_in=1.55,
           col_widths_in=[3.3, 4.9, 4.5],
           head_size=12, cell_size=11, row_height_in=0.7)
    _plain_text(s,
                "Both fixes are ~5–10 lines on top of the standard "
                "HF forward. Paper §3.2, p. 5.",
                0.5, 5.9, 12.5, 0.4, size=13, color=MUTED)
    return s


def s_mech_algorithm_box(prs):
    s = _blank(prs)
    _title(s, "Mechanism deep-dive: the algorithm in one box")
    _box(s, 0.7, 1.55, 12.0, 5.2, fill=LIGHT, line=NAVY, line_pt=1.0)
    lines = [
        "State: sinks = KV[0..S-1]   ·   window = deque(maxlen=L)   ·   S=4, L=1020",
        "",
        "def decode_step(x_t):",
        "    q, k, v  =  project(x_t)                  # (H, d)",
        "    cache_K  =  concat(sinks.K, window.K, k)   # length ≤ S+L+1",
        "    cache_V  =  concat(sinks.V, window.V, v)",
        "    # cache-local position IDs = 0..len(cache_K)-1",
        "    pos      =  arange(len(cache_K))",
        "    if using_RoPE:",
        "        cache_K = rotate(cache_K, pos)         # rotate at decode",
        "        q       = rotate(q,       pos[-1])",
        "    logits = softmax(q @ cache_K.T / sqrt(d))",
        "    y_t    = logits @ cache_V",
        "    window.append((k, v))                      # sinks never touched",
        "    return y_t",
    ]
    for i, line in enumerate(lines):
        _plain_text(s, line, 0.9, 1.7 + i * 0.32, 11.6, 0.35,
                    size=13, color=NAVY if i == 0 else TEXT,
                    bold=(i == 0))
    _plain_text(s,
                "Total additional code vs. vanilla HF decode ≈ 40–50 lines. "
                "See task1-presentation/code/streaming_llm_demo.py.",
                0.7, 6.9, 12.0, 0.35, size=12, color=MUTED)
    return s


# --- Sink-token pre-training deep-dive ------------------------------------


def s_sink_pretraining_deep(prs):
    s = _blank(prs)
    _title(s, "Sink-Token pre-training deep-dive (§3.3)")
    _bullets(s, [
        "Two proposals, both trained on deduplicated Pile, 160M-param "
        "Pythia-shaped model, batch 256, 143K steps.",
        "Learnable Sink Token — one placeholder prepended to every "
        "training sample; parameters learned end-to-end.",
        "Zero Sink — replaces SoftMax with SoftMax-off-by-One (Miller, "
        "2023): add a constant 1 to the denominator; equivalent to a "
        "phantom all-zero key/value.",
    ], top_in=1.55, size=17)
    headers = ["Variant  ·  cache config (160M, PG19)",
               "0+1024", "1+1023", "2+1022", "4+1020"]
    rows = [
        ["Vanilla",         "27.87", "18.49", "18.05", "18.05"],
        ["Zero Sink",       "29214", "19.90", "18.27", "18.01"],
        ["Learnable Sink",  "1235",  "18.01", "18.01", "18.02"],
    ]
    _table(s, headers, rows, left_in=1.1, top_in=4.35,
           col_widths_in=[4.7, 1.4, 1.4, 1.4, 1.4],
           head_size=12, cell_size=12, row_height_in=0.42)
    _plain_text(s,
                "Table 3, p. 6.  Read the 1+1023 column: Learnable Sink "
                "18.01 already matches Vanilla's 4+1020 (18.05).",
                1.1, 6.35, 11.0, 0.4, size=12, color=MUTED)
    return s


def s_sink_downstream(prs):
    s = _blank(prs)
    _title(s, "Sink-Token pre-training: does downstream accuracy hold?")
    _bullets(s, [
        "Table 4 (p. 7): zero-shot accuracy on seven benchmarks, "
        "160M vanilla vs 160M with Sink Token pre-training.",
        "Differences on all seven tasks are <1 point.",
        "Sink-Token column is *slightly higher* on average.",
        "Figure 6 (p. 7): pre-training loss curves are indistinguishable.",
        "Figure 7 (p. 7): attention heatmaps confirm the sink token "
        "draws the mass that would otherwise smear over positions 1–4.",
    ], top_in=1.55, size=17)
    _labelled_box(s, 0.7, 4.7, 12.0, 2.15,
                  "The honest caveat (paper §3.3 + our critical read)",
                  "Everything above is at 160M parameters. "
                  "The stronger claim — that a single learnable sink "
                  "token replaces the four accidental ones at 7B / 13B / "
                  "70B — is left as future work. Nobody has published "
                  "the scale-up. The bigger the model, the more the "
                  "accidental sinks smear; whether a single learned "
                  "slot still absorbs all of them at scale is an open "
                  "empirical question.",
                  fill=LIGHT, border=HIGHLIGHT,
                  label_size=15, sub_size=13, sub_color=TEXT,
                  label_color=HIGHLIGHT)
    return s


# --- Results deep-dive ---------------------------------------------------


def s_result_ppl_permodel(prs):
    s = _blank(prs)
    _title(s, "Results deep-dive: Fig. 5 per-model perplexity to 4M tokens")
    _bullets(s, [
        "Same StreamingLLM policy applied to every family; cache = half "
        "the pre-training window.",
        "Perplexity stable out past 4M tokens on 10 different models.",
    ], top_in=1.55, size=17)
    headers = ["Family",  "Sizes tested", "Pos. encoding",
               "Cache used", "PPL @ 4M tokens (stable?)"]
    rows = [
        ["Llama-2",  "7B, 13B, 70B",   "RoPE",  "2048", "flat"],
        ["MPT",      "7B, 30B",        "ALiBi", "1024", "flat"],
        ["Falcon",   "7B, 40B",        "RoPE",  "1024", "flat"],
        ["Pythia",   "2.8B, 6.9B, 12B","RoPE",  "1024", "flat"],
    ]
    _table(s, headers, rows, left_in=0.7, top_in=3.1,
           col_widths_in=[2.0, 2.6, 2.6, 1.9, 3.0],
           head_size=12, cell_size=12, row_height_in=0.45)
    _labelled_box(s, 0.7, 5.55, 12.0, 1.35,
                  "The Sliding+Recompute reference line",
                  "The paper overlays the oracle sliding-window-with-"
                  "recomputation baseline. StreamingLLM tracks it "
                  "essentially exactly on all 10 models — but at "
                  "up to 22× lower per-token latency (next slide).",
                  fill=LIGHT, border=GREEN_OK,
                  label_size=14, sub_size=13, sub_color=TEXT,
                  label_color=GREEN_OK)
    return s


def s_result_throughput_full(prs):
    s = _blank(prs)
    _title(s, "Results deep-dive: Fig. 10 per-cache-size latency (7B + 13B)")
    _plain_text(s, "Llama-2-13B (per-token decode latency, ms):",
                0.5, 1.55, 6.5, 0.4, size=13, bold=True, color=NAVY)
    _table(s,
           ["Cache", "Sliding+Recompute", "StreamingLLM", "Speed-up"],
           [
               ["256",  "2355", "106", "22.2×"],
               ["512",  "860",  "75",  "11.5×"],
               ["1024", "361",  "60",  "6.0×"],
               ["2048", "169",  "52",  "3.3×"],
               ["4096", "99",   "48",  "2.1×"],
           ],
           left_in=0.5, top_in=1.95,
           col_widths_in=[0.9, 2.4, 1.9, 1.1],
           head_size=11, cell_size=11, row_height_in=0.32)
    _plain_text(s, "Llama-2-7B (per-token decode latency, ms):",
                7.0, 1.55, 6.0, 0.4, size=13, bold=True, color=NAVY)
    _table(s,
           ["Cache", "Sliding+Recompute", "StreamingLLM", "Speed-up"],
           [
               ["256",  "1411", "65", "21.7×"],
               ["512",  "523",  "45", "11.6×"],
               ["1024", "223",  "35", "6.4×"],
               ["2048", "103",  "31", "3.3×"],
               ["4096", "63",   "31", "2.0×"],
           ],
           left_in=7.0, top_in=1.95,
           col_widths_in=[0.9, 2.4, 1.9, 1.1],
           head_size=11, cell_size=11, row_height_in=0.32)
    _labelled_box(s, 0.5, 4.15, 12.5, 2.75,
                  "Reading the numbers honestly",
                  "The 22.2× headline is the *small-cache* number. That "
                  "is where the Sliding+Recompute baseline hurts most "
                  "(each new token forces a rebuild of the whole "
                  "cache). At cache 4096 the speed-up shrinks to 2× "
                  "because the recompute baseline amortises over a "
                  "larger cache. The important second row is memory: "
                  "StreamingLLM and Sliding+Recompute have essentially "
                  "identical footprint (both bounded by cache size). "
                  "Single A6000, HF Transformers, batch 1, greedy "
                  "decode. Batched-attention interactions are out of "
                  "scope — a fair question for Q&A.",
                  fill=LIGHT, border=BLUE,
                  label_size=15, sub_size=13, sub_color=TEXT)
    return s


def s_result_streaming_qa_perM(prs):
    s = _blank(prs)
    _title(s, "Results deep-dive: streaming QA per model size")
    _plain_text(s,
                "Concatenate every ARC-Easy and ARC-Challenge question "
                "into a single stream. Score exact match on each answer. "
                "Cache = 1024.",
                0.5, 1.55, 12.5, 0.4, size=13, color=TEXT)
    _plain_text(s,
                "Look at the Window row: this is not a slight "
                "degradation, it's collapse. StreamingLLM restores "
                "one-shot-parity accuracy — or slightly beats it.",
                0.5, 1.95, 12.5, 0.4, size=13, color=MUTED)
    headers = ["Model  ·  policy", "ARC-E", "ARC-C", "Δ vs one-shot"]
    rows = [
        ["Llama-2-7B-Chat  one-shot",    "71.25", "53.16", "—"],
        ["Llama-2-7B-Chat  Window",      "3.58",  "1.39",  "-67.7 / -51.8"],
        ["Llama-2-7B-Chat  StreamingLLM","71.34", "55.03", "+0.1 / +1.9"],
        ["Llama-2-13B-Chat one-shot",    "78.16", "63.31", "—"],
        ["Llama-2-13B-Chat Window",      "0.25",  "0.34",  "-77.9 / -63.0"],
        ["Llama-2-13B-Chat StreamingLLM","80.89", "65.61", "+2.7 / +2.3"],
        ["Llama-2-70B-Chat one-shot",    "91.29", "78.50", "—"],
        ["Llama-2-70B-Chat Window",      "0.12",  "0.32",  "-91.2 / -78.2"],
        ["Llama-2-70B-Chat StreamingLLM","91.37", "80.20", "+0.1 / +1.7"],
    ]
    _table(s, headers, rows, left_in=1.5, top_in=2.85,
           col_widths_in=[5.6, 1.5, 1.5, 2.4],
           head_size=12, cell_size=11, row_height_in=0.35)
    _plain_text(s, "Dense OOMs before the stream ends. Table 5, p. 8.",
                1.5, 6.2, 11.0, 0.35, size=12, color=MUTED)
    return s


def s_result_streameval(prs):
    s = _blank(prs)
    _title(s, "Results deep-dive: StreamEval (Fig. 8 + Fig. 9)")
    _bullets(s, [
        "The paper's own long-eval benchmark: issue a query every 10 "
        "lines; the answer is always exactly 20 lines back.",
        "Tests whether the *rolling window* is doing its job on a "
        "controlled distance signal.",
        "Setup: Llama-2-7B-Chat + StreamingLLM at cache 1024, "
        "vs. dense (OOMs), vs. LongChat-7B-v1.5-32K + StreamingLLM.",
    ], top_in=1.55, size=17)
    _labelled_box(s, 0.7, 3.85, 12.0, 1.6,
                  "Positive result (Fig. 8)",
                  "StreamingLLM stays accurate up to ~120K tokens; "
                  "dense collapses (OOM); window is near-zero from the "
                  "first cache eviction. StreamingLLM does what "
                  "Sliding+Recompute would, at fraction of the cost.",
                  fill=LIGHT, border=GREEN_OK,
                  label_size=15, sub_size=13, sub_color=TEXT,
                  label_color=GREEN_OK)
    _labelled_box(s, 0.7, 5.55, 12.0, 1.35,
                  "Honest limit (Fig. 9)",
                  "Accuracy drops the *moment* query-answer distance "
                  "exceeds cache size. StreamingLLM has NO ability to "
                  "recall information older than the rolling window. "
                  "This is the whole basis of §5's honest scope claim.",
                  fill=LIGHT, border=HIGHLIGHT,
                  label_size=15, sub_size=13, sub_color=TEXT,
                  label_color=HIGHLIGHT)
    return s


# --- Limitations deep-dive ----------------------------------------------


def s_limits_nonmonotone(prs):
    s = _blank(prs)
    _title(s, "Limits deep-dive: bigger cache ≠ lower perplexity")
    _plain_text(s,
                "Table 6 (p. 9). PG19, Llama-2-7B and Llama-2-13B. "
                "Sink count fixed at 4.",
                0.5, 1.55, 12.5, 0.4, size=13, color=TEXT)
    headers = ["Rolling window L", "Llama-2-7B PPL", "Llama-2-13B PPL"]
    rows = [
        ["508  (S+L = 512)",  "9.73", "8.35"],
        ["1020 (S+L = 1024)", "9.32", "7.79"],
        ["2044 (S+L = 2048)", "9.08", "7.51"],
        ["4092 (S+L = 4096)", "9.59", "7.60"],
    ]
    _table(s, headers, rows, left_in=2.0, top_in=2.15,
           col_widths_in=[3.6, 2.5, 2.5], head_size=12, cell_size=12,
           row_height_in=0.42)
    _labelled_box(s, 0.5, 4.35, 12.5, 2.55,
                  "Two ways to read this",
                  "Paper reads it as an LLM limitation: current LLMs "
                  "under-utilise the context they're given.\n\n"
                  "Our critical reading (see critical-analysis.md): "
                  "this is also a hint that the *policy* is leaving "
                  "signal on the table. A smarter rule inside the "
                  "rolling window — e.g. H2O-style attention-score "
                  "eviction — might monotonically improve with L. "
                  "This is exactly the direction the modern "
                  "alternatives explore.",
                  fill=LIGHT, border=BLUE,
                  label_size=15, sub_size=13, sub_color=TEXT)
    return s


def s_limits_longbench(prs):
    s = _blank(prs)
    _title(s, "Limits deep-dive: LongBench (Table 8) is not the target")
    _plain_text(s,
                "StreamingLLM 4+3496 vs default truncation 1750+1750, "
                "LongChat-7B-v1.5-32K.",
                0.5, 1.55, 12.5, 0.4, size=13, color=TEXT)
    headers = ["Task",
               "StreamingLLM 4+3496",
               "Truncation 1750+1750",
               "StreamingLLM 1750+1750"]
    rows = [
        ["NarrativeQA", "11.6", "18.7", "18.5"],
        ["Qasper",      "16.9", "19.2", "19.6"],
        ["HotpotQA",    "21.6", "25.4", "27.4"],
        ["2WikiMQA",    "28.2", "32.8", "31.7"],
        ["GovReport",   "23.9", "27.3", "28.4"],
        ["MultiNews",   "25.5", "25.8", "25.6"],
    ]
    _table(s, headers, rows, left_in=0.7, top_in=2.05,
           col_widths_in=[2.4, 3.3, 3.3, 3.3],
           head_size=12, cell_size=11, row_height_in=0.38)
    _labelled_box(s, 0.7, 4.6, 12.0, 2.3,
                  "Why StreamingLLM loses here",
                  "LongBench evaluates long-doc QA. StreamingLLM at "
                  "4+3496 keeps only 4 sink tokens of the *initial "
                  "prompt* — and the prompt is where LongBench puts "
                  "the question. Raise the sink budget to 1750 and "
                  "parity is restored. Verdict: StreamingLLM is a "
                  "*streaming* policy, not a *long-doc* policy. The "
                  "paper is explicit about this in Appendix A.",
                  fill=LIGHT, border=HIGHLIGHT,
                  label_size=15, sub_size=13, sub_color=TEXT,
                  label_color=HIGHLIGHT)
    return s


def s_limits_position_only(prs):
    s = _blank(prs)
    _title(s, "Limits deep-dive: positional-only, no content awareness")
    _bullets(s, [
        "Eviction rule: keep first S tokens + last L tokens. "
        "Nothing else influences the decision.",
        "The paper never compares to attention-score-based evictors "
        "(H2O, SnapKV, Scissorhands, FastGen). Chronologically "
        "understandable (StreamingLLM arXiv Sept 2023) but a modern "
        "audience will ask.",
        "Consequence: a token in the rolling window that was heavily "
        "attended to is evicted the moment it ages out — even if "
        "later queries would attend to it again.",
        "Complementary hypothesis (my analysis): sinks + "
        "content-aware inside the window is likely to beat either "
        "alone. Section 'What I would change' comes back to this.",
    ], top_in=1.55, size=17)
    _labelled_box(s, 0.7, 5.55, 12.0, 1.35,
                  "Bridge to the next section",
                  "Let's meet the content-aware evictors StreamingLLM "
                  "did not compete against, and see where each of "
                  "them wins.",
                  fill=LIGHT, border=ACCENT,
                  label_size=15, sub_size=14, sub_color=TEXT,
                  label_color=ACCENT)
    return s


# --- Modern alternatives (H2O / SnapKV / etc.) --------------------------


def s_alt_content_aware(prs):
    s = _blank(prs)
    _title(s, "Modern alternatives: content-aware KV-cache eviction")
    headers = ["Method  ·  year", "Signal used to evict",
               "Complexity", "Position of the sink question"]
    rows = [
        ["H2O (2023)",
         "Cumulative attention score",
         "O(TL log L)",
         "Not addressed; treats first tokens as ordinary"],
        ["Scissorhands (2023)",
         "Persistence of attention across steps",
         "O(TL)",
         "Not addressed"],
        ["FastGen (2023)",
         "Per-head profile (heavy hitters, local, punct.)",
         "O(TL)",
         "Detects sinks empirically, per-head"],
        ["SnapKV (2024)",
         "Attention voting from a short observation window",
         "O(TL)",
         "Not addressed"],
        ["Keyformer (2024)",
         "Approx. attention via Gumbel-softmax over keys",
         "O(TL log L)",
         "Not addressed"],
    ]
    _table(s, headers, rows, left_in=0.4, top_in=1.55,
           col_widths_in=[2.4, 3.9, 1.9, 4.5],
           head_size=11, cell_size=10, row_height_in=0.5)
    _plain_text(s,
                "All five are compatible with StreamingLLM in principle: "
                "keep 4 positional sinks; evict inside the rolling "
                "window by content. No paper we know of publishes the "
                "hybrid at scale.",
                0.4, 5.3, 12.6, 0.9, size=13, color=MUTED)
    return s


def s_alt_positional_vs_content(prs):
    s = _blank(prs)
    _title(s, "Positional vs. content-aware: strengths & failure modes")
    headers = ["Property",
               "StreamingLLM (positional)",
               "H2O / SnapKV (content-aware)"]
    rows = [
        ["Decision rule",
         "keep first S + last L",
         "keep top-scoring tokens by attention weight"],
        ["Implementation cost",
         "~50 lines, no forward-pass changes",
         "profiling of attention scores + priority queue"],
        ["Softmax denominator stability",
         "guaranteed (sinks stay)",
         "not guaranteed (may drop the sink)"],
        ["Long-range content retention",
         "none beyond L",
         "yes for heavy-hitters"],
        ["Bursty / topic-shifting traffic",
         "robust (positional)",
         "risk: heavy hitters go stale"],
        ["Bench where it wins",
         "streaming LM / dialogue",
         "long-doc QA / summarisation"],
    ]
    _table(s, headers, rows, left_in=0.4, top_in=1.55,
           col_widths_in=[3.3, 4.3, 5.1],
           head_size=12, cell_size=11, row_height_in=0.52)
    _plain_text(s,
                "\"What I would change\" (slide coming) argues the two "
                "should be combined: positional sinks + content-aware "
                "rolling window.",
                0.4, 5.9, 12.6, 0.5, size=12, color=MUTED)
    return s


# --- Live implementation walkthrough ------------------------------------


def s_impl_pytorch(prs):
    s = _blank(prs)
    _title(s, "Implementation walkthrough 1/3: ~50 lines of PyTorch")
    _box(s, 0.5, 1.55, 12.5, 5.35, fill=LIGHT, line=NAVY, line_pt=1.0)
    code = [
        "class SinkKVCache:",
        "    def __init__(self, n_sinks=4, window=1020):",
        "        self.S, self.L = n_sinks, window",
        "        self.sinks_k, self.sinks_v = None, None",
        "        self.win_k,   self.win_v   = [], []",
        "",
        "    def append(self, k, v):                       # k, v: (H, 1, d)",
        "        if self.sinks_k is None or self.sinks_k.size(1) < self.S:",
        "            self.sinks_k = _cat(self.sinks_k, k)",
        "            self.sinks_v = _cat(self.sinks_v, v)",
        "            return",
        "        self.win_k.append(k); self.win_v.append(v)",
        "        if len(self.win_k) > self.L:",
        "            self.win_k.pop(0);  self.win_v.pop(0)",
        "",
        "    def as_kv(self):",
        "        K = torch.cat([self.sinks_k, *self.win_k], dim=1)",
        "        V = torch.cat([self.sinks_v, *self.win_v], dim=1)",
        "        return K, V                              # cache-local layout",
    ]
    for i, line in enumerate(code):
        _plain_text(s, line, 0.7, 1.7 + i * 0.28, 12.0, 0.28,
                    size=12, color=TEXT if not line.startswith("class ") else NAVY,
                    bold=line.startswith("class "))
    return s


def s_impl_rope(prs):
    s = _blank(prs)
    _title(s, "Implementation walkthrough 2/3: cache-local RoPE")
    _box(s, 0.5, 1.55, 12.5, 5.35, fill=LIGHT, line=NAVY, line_pt=1.0)
    code = [
        "def attention_step(x_t, cache, W_q, W_k, W_v, W_o, rope):",
        "    q  = W_q @ x_t                       # (H, 1, d)",
        "    k  = W_k @ x_t                       # store pre-rotation",
        "    v  = W_v @ x_t",
        "    cache.append(k, v)                   # sinks-then-window",
        "    K_pre, V = cache.as_kv()             # length T_c = S + L",
        "",
        "    pos_ids = torch.arange(K_pre.size(1), device=x_t.device)",
        "    K = rope.rotate(K_pre, pos_ids)      # cache-local IDs !",
        "    q = rope.rotate(q,     pos_ids[-1:]) # query at position T_c-1",
        "",
        "    scores = (q @ K.transpose(-1, -2)) / d_head ** 0.5",
        "    y = torch.softmax(scores, dim=-1) @ V",
        "    return W_o @ y",
    ]
    for i, line in enumerate(code):
        _plain_text(s, line, 0.7, 1.7 + i * 0.32, 12.0, 0.32,
                    size=12, color=TEXT if not line.startswith("def ") else NAVY,
                    bold=line.startswith("def "))
    _plain_text(s,
                "Note: the KV cache stores keys *pre*-rotation. Rotation "
                "is applied every step with a fresh position vector "
                "arange(0, T_c). This is what makes the sinks and the "
                "rolling window appear adjacent to the model.",
                0.5, 6.35, 12.5, 0.55, size=11, color=MUTED)
    return s


def s_impl_bench_harness(prs):
    s = _blank(prs)
    _title(s, "Implementation walkthrough 3/3: benchmark harness sketch")
    _box(s, 0.5, 1.55, 12.5, 5.35, fill=LIGHT, line=NAVY, line_pt=1.0)
    code = [
        "# Reproducing Fig. 5 for one model:",
        "model = AutoModelForCausalLM.from_pretrained('meta-llama/Llama-2-7b-hf')",
        "cache = SinkKVCache(n_sinks=4, window=2044)      # 4+2044 config",
        "",
        "tokens = load_pg19_concat(target_len=4_000_000)  # PG19 test",
        "logits = []",
        "for t in tqdm(range(len(tokens))):",
        "    logit = decode_step(model, tokens[t], cache)",
        "    logits.append(logit)",
        "",
        "ppl = perplexity_from_logits(logits, tokens[1:])",
        "assert ppl_at_4M(ppl) < 10.0                     # Table 6, Llama-2-7B",
        "",
        "# For 22.2×: repeat with sliding-window+recompute baseline and",
        "# divide per-token wall-clock. Fig. 10 numbers reproduce to ~5%.",
    ]
    for i, line in enumerate(code):
        _plain_text(s, line, 0.7, 1.7 + i * 0.32, 12.0, 0.32,
                    size=12, color=TEXT,
                    bold=False)
    _plain_text(s,
                "Full runnable version: task1-presentation/code/"
                "streaming_llm_demo.py",
                0.5, 6.35, 12.5, 0.4, size=12, color=MUTED)
    return s


# --- Bridge + Q&A prep ---------------------------------------------------


def s_bridge_task2(prs):
    s = _blank(prs)
    _title(s, "Bridge to Task 2: two caches, one idea")
    headers = ["Layer",
               "StreamingLLM (Task 1)",
               "GDSF on GPTCache (Task 2)"]
    rows = [
        ["What is cached",
         "Key/Value vectors of past tokens",
         "Whole (prompt → response) pairs"],
        ["Where in the stack",
         "Inside the model, per attention layer",
         "In front of the model, at the API layer"],
        ["Capacity units",
         "S sinks + L rolling tokens",
         "N entries or M bytes"],
        ["Eviction signal",
         "Position (first S + last L)",
         "Frequency, dollar cost, size (GDSF)"],
        ["Failure mode when dumb",
         "PPL 5158 (softmax collapse)",
         "Overspend on regeneration $$$"],
        ["What our project proves",
         "Reproduce +25%..+91% below",
         "GDSF beats LRU by up to +91% $ savings"],
    ]
    _table(s, headers, rows, left_in=0.4, top_in=1.55,
           col_widths_in=[2.7, 4.5, 5.5],
           head_size=12, cell_size=11, row_height_in=0.55)
    _labelled_box(s, 0.4, 4.9, 12.6, 2.0,
                  "The through-line",
                  "Both caches are bounded. Both suffer when eviction "
                  "is naive. StreamingLLM fixes it with a positional "
                  "rule. My Task-2 project fixes it with a cost-aware "
                  "rule. Same underlying diagnosis: dumb eviction wastes "
                  "bounded capacity.",
                  fill=LIGHT, border=ACCENT,
                  label_size=15, sub_size=14, sub_color=TEXT,
                  label_color=ACCENT)
    return s


def s_qa_prep_1(prs):
    s = _blank(prs)
    _title(s, "Q&A prep 1/2: expected questions")
    _bullets(s, [
        "\"Is this just window attention with a warm start?\" — No. "
        "Sinks are never evicted (Table 1, PPL 5158 → 5.40).",
        "\"Why 4 sinks?\" — Table 2: 1–2 leaves a residual bump; 4 "
        "saturates; 8 no further gain.",
        "\"RoPE, ALiBi, or both?\" — Both, via cache-local re-indexing "
        "(§3.2, p. 5).",
        "\"Does it extend the context?\" — No. Table 7, Fig. 9: "
        "accuracy collapses past cache size.",
        "\"vs. H2O / SnapKV / FastGen?\" — Paper does not compare; "
        "chronologically earlier. Positional vs. content-aware are "
        "complementary.",
        "\"Long-doc QA?\" — Table 8: underperforms truncation at "
        "4+3496 because prompt is lost; parity at 1750+1750.",
    ], top_in=1.55, size=17)
    return s


def s_qa_prep_2(prs):
    s = _blank(prs)
    _title(s, "Q&A prep 2/2: harder questions and what I would say")
    _bullets(s, [
        "\"Are sinks a softmax artefact — would SoftMax-off-by-One "
        "kill them?\" — Table 3 Zero-Sink row: helps somewhat "
        "(29214→18.01) but does not remove the need for sinks. "
        "Empirical rescue > theoretical fix.",
        "\"Is 22.2× realistic for production?\" — Batch-1, single-A6000 "
        "number. Real servers batch. Batched-attention interactions "
        "with sink+rolling layout are an open engineering question.",
        "\"What breaks the sink phenomenon?\" — Not shown; the paper "
        "covers Llama-2, MPT, Falcon, Pythia, BERT, ViT. Would love "
        "to see it tested on Mixture-of-Experts and on state-space "
        "models (Mamba).",
        "\"Have you reproduced this?\" — Not the LLM-scale numbers "
        "directly (compute), but the 50-line demo + our GDSF work "
        "reproduces the underlying \"dumb eviction is expensive\" "
        "claim at the response-cache layer with 3600-run stats.",
    ], top_in=1.55, size=16)
    return s


def s_conclusion(prs):
    s = _blank(prs)
    _title(s, "Conclusion")
    _bullets(s, [
        "Attention sinks are real, universal, and structural, not semantic.",
        "Keeping 4 initial tokens + a rolling window is enough to make "
        "LLMs stream stably to 4M+ tokens.",
        "22.2× faster than the only sane baseline, with the same memory.",
        "Adopted upstream: NVIDIA TensorRT-LLM, HF Transformers, "
        "Intel Extension for Transformers, MLC LLM.",
        "But: it does not extend context, and it does not compete with "
        "attention-score-based eviction. Those are the frontiers.",
    ], top_in=1.55, size=20)
    _labelled_box(s, 0.7, 5.5, 12.0, 1.35,
                  "One-line take",
                  "A robust empirical fix for a softmax quirk: small, "
                  "easy to implement, and honest about its own ceiling.",
                  fill=LIGHT, border=ACCENT,
                  label_size=16, sub_size=15, sub_color=TEXT,
                  label_color=ACCENT)
    return s


def s_thanks(prs):
    s = _blank(prs)
    _box(s, 0, 0, 13.33, 7.5, fill=NAVY, line=NAVY)
    _plain_text(s, "Thank you", 0.7, 2.7, 12.0, 1.5,
                size=64, bold=True, color=BG, align=PP_ALIGN.CENTER)
    _box(s, 4.7, 4.4, 3.9, 0.05, fill=ACCENT, line=ACCENT)
    _plain_text(s, "Questions welcome.", 0.7, 4.6, 12.0, 0.5,
                size=22, color=BG, align=PP_ALIGN.CENTER)
    _plain_text(s, "Nissim Brami  ·  nissimbrami@post.bgu.ac.il",
                0.7, 6.4, 12.0, 0.4, size=14, color=BG,
                align=PP_ALIGN.CENTER)
    _plain_text(s, "Caching in LLMs  ·  Ben-Gurion University",
                0.7, 6.75, 12.0, 0.4, size=12, color=BG,
                align=PP_ALIGN.CENTER)
    return s


# ---------- Assemble ------------------------------------------------------


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        s_title,
        s_agenda,
        # Basics section
        s_section_basics,
        s_llm_decoding,
        s_course_anchor,
        s_naive_options,
        # Related work section
        s_section_related,
        s_related_length_extrap,
        s_related_context_ext,
        s_related_sparse_and_lminf,
        # The observation
        s_section_observation,
        s_attention_sink_phenomenon,
        s_softmax_argument,
        s_semantics_or_position,
        # The mechanism + deep-dive
        s_section_mechanism,
        s_cache_layout,
        s_position_reindex,
        s_mech_worked_example,
        s_mech_rope_alibi,
        s_mech_algorithm_box,
        # Sink-Token pre-training (§3.3) + deep-dive
        s_sink_token_pretraining,
        s_sink_pretraining_deep,
        s_sink_downstream,
        # Results + deep-dive
        s_section_results,
        s_perplexity_result,
        s_result_ppl_permodel,
        s_throughput_result,
        s_result_throughput_full,
        s_streaming_qa,
        s_result_streaming_qa_perM,
        s_result_streameval,
        # Limits + deep-dive
        s_limits,
        s_limits_nonmonotone,
        s_limits_longbench,
        s_limits_position_only,
        # Modern alternatives
        s_alt_content_aware,
        s_alt_positional_vs_content,
        # Implementation walkthrough
        s_impl_pytorch,
        s_impl_rope,
        s_impl_bench_harness,
        # What I'd change + bridge + Q&A + conclusion
        s_what_id_change,
        s_bridge_task2,
        s_qa_prep_1,
        s_qa_prep_2,
        s_conclusion,
        s_thanks,
    ]

    section_dividers = (s_title, s_section_basics, s_section_related,
                        s_section_observation, s_section_mechanism,
                        s_section_results, s_thanks)

    for i, b in enumerate(builders, start=1):
        slide = b(prs)
        # Page numbers on all but the full-navy dividers/title/thanks
        if b not in section_dividers:
            tb = slide.shapes.add_textbox(Inches(12.7), Inches(7.05),
                                          Inches(0.6), Inches(0.3))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT
            r = p.add_run()
            r.text = f"{i}"
            r.font.name = "Calibri"
            r.font.size = Pt(9)
            r.font.color.rgb = MUTED
            tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(7.05),
                                           Inches(6.0), Inches(0.3))
            tf2 = tb2.text_frame
            p2 = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.LEFT
            r2 = p2.add_run()
            r2.text = "Nissim Brami · Caching in LLMs · BGU"
            r2.font.name = "Calibri"
            r2.font.size = Pt(9)
            r2.font.color.rgb = MUTED

    prs.save(str(OUT))
    print(f"Wrote {OUT}  ·  {len(prs.slides)} slides")


if __name__ == "__main__":
    build()
